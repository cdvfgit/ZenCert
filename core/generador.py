"""
generador.py
------------
Motor de generación de certificados PPTX.

Recibe datos ya procesados por procesador.py y una ruta de plantilla
válida, y produce el archivo PPTX final. No formatea datos ni llama
a APIs externas — su única responsabilidad es escribir el PPTX.

Método: python-pptx reemplaza únicamente el texto de los placeholders.
Fuentes, tamaños, colores, posición y centrado quedan 100% intactos
tal como están definidos en el PPTX original.

PLACEHOLDERS ESPERADOS EN LA PLANTILLA:
    {{Nombre_Alumno}}     → nombre completo del alumno
    {{cedula}}            → cédula formateada (o vacío si no aplica)
    {{grado_español}}     → grado en español con color de cinturón
    {{ciudad_fecha}}      → ciudad y fecha del examen
    {{registro}}          → número de registro formateado
    {{nombre_katakana}}   → nombre transliterado en katakana
    {k}                   → kanji del grado
    {{fecha_japonesa}}    → fecha completa en japonés
"""

import copy
import json
import logging
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.opc.packuri import PackURI
from pptx.parts.image import ImagePart

logger = logging.getLogger(__name__)

# Namespace de relaciones para reescribir r:embed y r:link en el XML
_R_NS  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_EMBED = f"{{{_R_NS}}}embed"
_LINK  = f"{{{_R_NS}}}link"

_RUTA_ORGANIZACIONES = Path(__file__).parent.parent / "data" / "organizaciones"


def _cargar_placeholders(organizacion: str) -> dict:
    """
    Carga el mapa de placeholders desde data/organizaciones/{org}.json.

    Args:
        organizacion: Codigo de la organizacion (ej. 'ocoa').

    Returns:
        Diccionario placeholder a clave de datos.

    Raises:
        FileNotFoundError: Si no existe el archivo de la organizacion.
        KeyError: Si la organizacion no tiene placeholders definidos.
    """
    ruta = _RUTA_ORGANIZACIONES / f"{organizacion}.json"

    if not ruta.exists():
        raise FileNotFoundError(
            f"No se encontro configuracion para '{organizacion}'. "
            f"Ruta esperada: {ruta}. "
            "Verifica que exista el archivo en data/organizaciones/"
        )

    config = json.loads(ruta.read_text(encoding="utf-8"))
    placeholders = config.get("placeholders")

    if not placeholders:
        raise KeyError(
            f"No hay placeholders definidos para '{organizacion}'. "
            f"Agrega la clave 'placeholders' en {ruta.name}"
        )

    return placeholders


# ── Helpers privados ──────────────────────────────────────────────────────────

def _reemplazar_en_shape(shape, datos: dict, placeholders: dict) -> None:
    """
    Reemplaza placeholders en un shape conservando todo el formato original.

    Maneja dos casos:
    - Placeholder en un solo run: reemplazo directo conservando formato.
    - Placeholder dividido en múltiples runs por PowerPoint: reconstruye
      el texto completo del párrafo, hace el reemplazo, escribe el resultado
      en el primer run y limpia los runs siguientes. El formato del primer
      run se conserva intacto.

    Este comportamiento es transparente — no afecta shapes sin placeholders
    ni plantillas donde los placeholders ya están en un solo run.

    Args:
        shape: Shape de python-pptx a procesar.
        datos: Diccionario con los valores a sustituir.
        placeholders: Mapa placeholder → clave cargado desde organizaciones/{org}.json.
    """
    if not shape.has_text_frame:
        return

    for para in shape.text_frame.paragraphs:
        if not para.runs:
            continue

        # Reconstruir texto completo del párrafo uniendo todos los runs
        texto_completo = "".join(r.text for r in para.runs)

        # Verificar si hay algún placeholder en el texto completo
        texto_reemplazado = texto_completo
        hay_reemplazo = False

        for placeholder, clave in placeholders.items():
            if placeholder in texto_reemplazado:
                texto_reemplazado = texto_reemplazado.replace(
                    placeholder,
                    datos.get(clave, ""),
                )
                hay_reemplazo = True

        if not hay_reemplazo:
            continue

        # Escribir el texto reemplazado en el primer run y limpiar los demás
        # El formato (fuente, tamaño, color) del primer run se conserva intacto
        para.runs[0].text = texto_reemplazado
        for run in para.runs[1:]:
            run.text = ""


def _clonar_slide(prs_destino: Presentation, slide_origen) -> None:
    """
    Clona slide_origen al final de prs_destino con ImageParts únicas.

    Cada imagen se registra como una ImagePart nueva con nombre de
    archivo único dentro del ZIP, evitando el error:
        UserWarning: Duplicate name: 'ppt/media/imageX.png'

    Args:
        prs_destino: Presentación destino donde se agrega la slide.
        slide_origen: Slide ya procesada con datos reemplazados.
    """
    # Contar imágenes existentes para generar índices únicos
    n_imgs = sum(
        1 for s in prs_destino.slides
        for rel in s.part.rels.values()
        if "image" in rel.reltype
    )

    nueva_slide = prs_destino.slides.add_slide(prs_destino.slide_layouts[0])

    # Copiar árbol XML de shapes desde la slide origen
    nueva_slide.shapes._spTree.clear()
    for elem in slide_origen.shapes._spTree:
        nueva_slide.shapes._spTree.append(copy.deepcopy(elem))

    # Crear ImageParts nuevas y únicas para cada imagen de la slide origen
    rId_map: dict[str, str] = {}
    for rId_orig, rel in slide_origen.part.rels.items():
        if "image" not in rel.reltype:
            continue

        img_orig = rel._target
        n_imgs += 1
        ext = "png" if "png" in img_orig.content_type else "jpg"

        nuevo_partname = PackURI(
            f"/ppt/media/img_{len(prs_destino.slides)}_{n_imgs}.{ext}"
        )

        nueva_img = ImagePart(
            nuevo_partname,
            img_orig.content_type,
            prs_destino.part.package,
            img_orig.blob,
        )

        nuevo_rId = nueva_slide.part.relate_to(nueva_img, rel.reltype)
        rId_map[rId_orig] = nuevo_rId

    # Reescribir referencias r:embed y r:link para apuntar a los nuevos rIds
    for elem in nueva_slide.shapes._spTree.iter():
        for attr in (_EMBED, _LINK):
            old_rid = elem.get(attr)
            if old_rid and old_rid in rId_map:
                elem.set(attr, rId_map[old_rid])


def _ruta_salida(id_orden: str, organizacion: str) -> Path:
    """
    Construye y garantiza la existencia de la carpeta de salida para la orden.

    Args:
        id_orden: Identificador de la orden (ej. 'ORD-2025-001').
        organizacion: Código de la organización (ej. 'ocoa').

    Returns:
        Path de la carpeta de salida lista para escribir.
    """
    ruta = Path(__file__).parent.parent / "salida" / organizacion / id_orden
    ruta.mkdir(parents=True, exist_ok=True)
    return ruta


# ── API pública ───────────────────────────────────────────────────────────────

def detectar_plantilla(
    organizacion: str,
    tipo: str,
    dojo: str,
) -> Path:
    """
    Detecta la plantilla PPTX correcta usando busqueda por prioridad.

    Busca en orden descendente de especificidad dentro de
    plantillas/{organizacion}/. Usa la primera que encuentre.

    Prioridad:
        1. {org}_{tipo}_{dojo}.pptx   mas especifica
        2. {org}_{dojo}.pptx          sin tipo
        3. {org}_{tipo}.pptx          sin dojo

    Args:
        organizacion: Codigo de la organizacion (ej. 'ocoa').
        tipo: Tipo de grado (ej. 'dan', 'kyu').
        dojo: Nombre normalizado del dojo (ej. 'puerto_ordaz').

    Returns:
        Path absoluto de la plantilla encontrada.

    Raises:
        FileNotFoundError: Si ninguna plantilla coincide con los datos
                           de la orden. El estado de la orden no cambia.
    """
    base = Path(__file__).parent.parent / "plantillas" / organizacion

    candidatos = [
        base / f"{organizacion}_{tipo}_{dojo}.pptx",
        base / f"{organizacion}_{dojo}.pptx",
        base / f"{organizacion}_{tipo}.pptx",
    ]

    for ruta in candidatos:
        if ruta.exists():
            logger.info(
                "detectar_plantilla | encontrada | org=%s | tipo=%s | dojo=%s | archivo=%s",
                organizacion, tipo, dojo, ruta.name,
            )
            return ruta

    raise FileNotFoundError(
        f"No se encontro plantilla para:\n"
        f"  Organizacion : {organizacion}\n"
        f"  Tipo         : {tipo}\n"
        f"  Dojo         : {dojo}\n"
        f"Archivos buscados:\n"
        + "\n".join(f"  - {c.name}" for c in candidatos)
        + f"\nVerifica que exista alguno en plantillas/{organizacion}/"
    )


def generar_lote(
    lista_alumnos: list[dict],
    ruta_plantilla: str | Path,
    id_orden: str,
    organizacion: str,
    nombre_salida: str = "certificados_lote.pptx",
) -> str:
    """
    Genera un único PPTX donde cada diapositiva es un certificado.

    Cada slide parte de una copia fresca del PPTX base, garantizando
    que los datos de un alumno no afecten a los demás. Las imágenes
    de cada slide se registran con nombres únicos en el ZIP.

    Args:
        lista_alumnos: Lista de dicts procesados por procesador.procesar_orden().
                       Cada dict debe tener las claves definidas en PLACEHOLDERS.
        ruta_plantilla: Ruta absoluta o relativa al PPTX base de la organización.
        id_orden: Identificador de la orden (ej. 'ORD-2025-001').
                  Se usa para construir la carpeta de salida.
        organizacion: Código de la organización (ej. 'ocoa').
                      Se usa para construir la carpeta de salida.
        nombre_salida: Nombre del archivo PPTX resultante
                       (default: 'certificados_lote.pptx').

    Returns:
        Ruta absoluta del archivo PPTX generado.

    Raises:
        FileNotFoundError: Si la plantilla PPTX no existe en la ruta indicada.
        ValueError: Si la lista de alumnos está vacía.
    """
    ruta_plantilla = Path(ruta_plantilla)

    if not ruta_plantilla.exists():
        raise FileNotFoundError(
            f"No se encontró la plantilla en: {ruta_plantilla}\n"
            "Verifica que el archivo exista en plantillas/ con el nombre correcto."
        )

    if not lista_alumnos:
        raise ValueError(
            f"La lista de alumnos para la orden '{id_orden}' está vacía."
        )

    logger.info(
        "generar_lote | id_orden=%s | plantilla=%s | total=%d",
        id_orden,
        ruta_plantilla.name,
        len(lista_alumnos),
    )

    # Cargar placeholders una sola vez para todo el lote
    placeholders = _cargar_placeholders(organizacion)

    prs = Presentation(str(ruta_plantilla))

    for i, datos in enumerate(lista_alumnos):
        nombre = datos.get("nombre_alumno", f"alumno_{i + 1}")

        if i == 0:
            # Primera slide: reemplazar directamente en la slide original
            for shape in prs.slides[0].shapes:
                _reemplazar_en_shape(shape, datos, placeholders)
        else:
            # Slides siguientes: copia fresca del base → reemplazar → clonar
            prs_tmp = Presentation(str(ruta_plantilla))
            slide_tmp = prs_tmp.slides[0]

            for shape in slide_tmp.shapes:
                _reemplazar_en_shape(shape, datos, placeholders)

            _clonar_slide(prs, slide_tmp)

        logger.debug("generar_lote | slide=%d | alumno=%s", i + 1, nombre)

    carpeta = _ruta_salida(id_orden, organizacion)
    ruta_archivo = carpeta / nombre_salida
    prs.save(str(ruta_archivo))

    logger.info("generar_lote | completado | ruta=%s", ruta_archivo.resolve())
    return str(ruta_archivo.resolve())


def generar_certificado(
    datos: dict,
    ruta_plantilla: str | Path,
    id_orden: str,
    organizacion: str,
    nombre_archivo: str | None = None,
) -> str:
    """
    Genera un certificado PPTX individual para un solo alumno.

    Wrapper de generar_lote() para el caso de un único alumno.

    Args:
        datos: Dict procesado por procesador.procesar_orden() para un alumno.
        ruta_plantilla: Ruta al PPTX base de la organización.
        id_orden: Identificador de la orden.
        organizacion: Código de la organización.
        nombre_archivo: Nombre del PPTX de salida sin extensión.
                        Si no se indica, usa el nombre del alumno.

    Returns:
        Ruta absoluta del archivo PPTX generado.

    Raises:
        FileNotFoundError: Si la plantilla no existe.
    """
    if not nombre_archivo:
        nombre_base = datos.get("nombre_alumno", "certificado")
        nombre_archivo = "".join(
            c for c in nombre_base if c.isalnum() or c in " ._-"
        ).strip()

    return generar_lote(
        lista_alumnos=[datos],
        ruta_plantilla=ruta_plantilla,
        id_orden=id_orden,
        organizacion=organizacion,
        nombre_salida=f"{nombre_archivo}.pptx",
    )

if __name__ == "__main__":

    from procesador import procesar_orden
    from registro import leer_ultimo, confirmar, reservar_rango

    orden = {
        "id_orden":     "ORD-2025-001",
        "organizacion": "ocoa",
        "dojo":         "puerto_ordaz",
        "ciudad":       "Caracas",
        "fecha":        date(2025, 2, 20),
    }


    alumnos = [
        { "nombre_alumno": "Juan Pérez",  "cedula": "25446976",
          "prefijo_cedula": "V", "grado": "1", "tipo": "dan", "notas": "" },
        { "nombre_alumno": "María Gómez", "cedula": "30123456",
          "prefijo_cedula": "V", "grado": "3", "tipo": "dan", "notas": "" },
        { "nombre_alumno": "Mario Pérez",  "cedula": "254",
          "prefijo_cedula": "E", "grado": "5", "tipo": "dan", "notas": "" },
        { "nombre_alumno": "María Gómez", "cedula": "",
          "prefijo_cedula": "V", "grado": "6", "tipo": "dan", "notas": "" },
        { "nombre_alumno": "Georgina Pérez",  "cedula": "254",
          "prefijo_cedula": "E", "grado": "10", "tipo": "dan", "notas": "" },
        { "nombre_alumno": "María rrrrrr", "cedula": "",
          "prefijo_cedula": "V", "grado": "6", "tipo": "dan", "notas": "" }
    ]

    grupos = procesar_orden(orden, alumnos, registro_inicio=138, modo_prueba=True)

    for tipo, alumnos_grupo in grupos.items():
        ruta_plantilla = detectar_plantilla(orden["organizacion"], tipo, orden["dojo"])
        ruta = generar_lote(
            lista_alumnos=alumnos_grupo,
            ruta_plantilla=ruta_plantilla,
            id_orden=orden["id_orden"],
            organizacion=orden["organizacion"],
            nombre_salida=f"certificados_{tipo}.pptx",
        )
        print(f"✅ {tipo}: {ruta}")
