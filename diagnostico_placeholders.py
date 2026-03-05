"""
diagnostico_placeholders.py
----------------------------
Muestra todos los textos encontrados en los shapes de una plantilla PPTX.
Útil para verificar que los placeholders en el archivo coinciden exactamente
con los definidos en data/organizaciones/{org}.json

Uso:
    python diagnostico_placeholders.py plantillas/nueva_org/plantilla.pptx
"""

import sys
from pptx import Presentation

def diagnosticar(ruta_pptx: str):
    prs = Presentation(ruta_pptx)
    slide = prs.slides[0]

    print(f"\n📄 Plantilla: {ruta_pptx}")
    print("─" * 60)
    print(f"{'SHAPE':<30} {'PÁRRAFO':<5} {'RUN':<5} TEXTO")
    print("─" * 60)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p_idx, para in enumerate(shape.text_frame.paragraphs):
            # Mostrar texto completo del párrafo
            texto_para = "".join(r.text for r in para.runs)
            if texto_para.strip():
                print(f"  {shape.name:<28} p{p_idx:<4} —     [{texto_para}]")

            # Mostrar cada run por separado para detectar splits
            for r_idx, run in enumerate(para.runs):
                if run.text.strip():
                    print(f"  {'':28} p{p_idx:<4} r{r_idx:<4} [{run.text}]")

    print("─" * 60)
    print("\n⚠  Si un placeholder aparece dividido en múltiples runs")
    print("   (ej: [{] [{{N] [ombre}] [}]) el reemplazo NO funcionará.")
    print("   Solución: en PowerPoint borra el texto y reescríbelo manualmente.\n")

if __name__ == "__main__":
    diagnosticar("plantillas/ogkv/ogkv_dan.pptx")