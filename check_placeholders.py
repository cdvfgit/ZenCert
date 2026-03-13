#!/usr/bin/env python3
"""
Script para extraer texto de PPTX y verificar placeholders
"""

import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent / "core"))

from pptx import Presentation

def extraer_texto_pptx(ruta_pptx):
    """Extrae todo el texto de un PPTX"""
    prs = Presentation(ruta_pptx)
    texto_completo = []
    
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texto_completo.append(f"Slide {slide_num + 1}: {shape.text.strip()}")
    
    return texto_completo

def main():
    # Revisar plantilla de ogkv_dan
    ruta = "plantillas/ogkv/ogkv_dan.pptx"
    
    if Path(ruta).exists():
        print(f"📄 Analizando: {ruta}")
        print("=" * 60)
        
        lineas = extraer_texto_pptx(ruta)
        for linea in lineas:
            print(linea)
            
            # Buscar placeholders
            if "{{" in linea and "}}" in linea:
                print(f"  🎯 PLACEHOLDER ENCONTRADO: {linea}")
    else:
        print(f"❌ No existe: {ruta}")

if __name__ == "__main__":
    main()
